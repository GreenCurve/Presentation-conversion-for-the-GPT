from pptx import Presentation
import pymupdf
import io
import os
import pytesseract
from PIL import Image

def folder_cr(folder_name,notice):
    '''
    Tries to create a folder; outputs message on finding it already exists
    '''
    if not os.path.exists(folder_name):
        os.makedirs(folder_name)
    else:
        print(notice)


def pptx_converter(presentation,out_name,middle_folder):
    '''
    Splice pptx file into list with entries corrseponding to the page number, all the text , paths to all images. 
    Creates all the images with reports for their text
    '''


    slide_text = [] #will contain all the text
    images_list = [] #will contain all the image paths

    for slide_num, slide in enumerate(presentation.slides, start=1): #iterate over the slides

        text = []
        images_in_page = []

        #iterate over all shapes in presentation
        for shape_index, shape in enumerate(slide.shapes): #for each element on slides

            if shape.has_text_frame: # if text - get it and write
                text.append(shape.text_frame.text)

            if shape.shape_type == 13:  # if image -  get it (13 corresponds to PICTURE type in pptx)
                image = shape.image 
                image_bytes = image.blob

                #define all paths
                image_filename = f"slide_{slide_num}_image_{shape_index}"
                folder_shape_name = '.' + '\\' + out_name + '\\' + middle_folder + '\\' + image_filename
                image_filename_full = out_name + "/" + middle_folder +  "/" + f"slide_{slide_num}_image_{shape_index}" + "/" + image_filename + ".png"


                folder_cr(folder_shape_name,'Directory for that image already exists!')


                with open(image_filename_full, "wb") as img_file:
                    img_file.write(image_bytes)


                image_report_filename = image_filename + ".txt"
                with open(out_name + "/" + middle_folder +  "/" + f"slide_{slide_num}_image_{shape_index}" + "/" + image_report_filename,"w",encoding="utf-8") as img_file_report:#create image in the Logs/
                    img = Image.open(image_filename_full)

                    # text_from_image = pytesseract.image_to_string(img)
                    # img_file_report.write(text_from_image)


                images_in_page.append(image_filename) #and write image filename

        images_list.append(images_in_page)
        slide_text.append(text)

    slides_content = []
    for slide_num, slide_content in enumerate(presentation.slides):
        # print(slide_content)
        slides_content.append({
            "slide_number": slide_num + 1,
            "text": slide_text[slide_num],
            "images": images_list[slide_num]
        })


    return slides_content



def pdf_converter(pdf_file,out_name,middle_folder):
    '''
    Splice pdf file into list with entries corrseponding to the page number, all the text , paths to all images. 
    Creates all the images with reports for their text
    '''

    pages_text = [] #will contain all the text
    pages_images = [] #will contain all the image paths

    for page_num in range(len(pdf_file)): #iterate over every page
        page = pdf_file[page_num]

        text = page.get_text("text") #get all text from page
        pages_text.append(text)#write all text


        image_list = page.get_images(full=True)  #get all images from page

        images_in_page = []

        #iterate over all images on the page
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_file.extract_image(xref)
            image_bytes = base_image["image"]

            # define paths
            image_filename = f"page_{page_num}_image_{img_index}"
            folder_image_name = '.' + '\\' + out_name + '\\' + middle_folder + '\\' + image_filename
            image_filename_full = out_name + "/" + middle_folder +  "/" + image_filename + "/" + image_filename + ".png"


            folder_cr(folder_image_name,'Directory for that image already exists!')


            with open(image_filename_full, "wb") as img_file:
                img_file.write(image_bytes)

            image_report_filename = image_filename + ".txt"
            with open(out_name + "/" + middle_folder +  "/" + image_filename + "/" + image_report_filename,"w",encoding="utf-8") as img_file_report:
                img = Image.open(image_filename_full)
                text_from_image = pytesseract.image_to_string(img, config="--oem 3 --psm 7")
                img_file_report.write(text_from_image)


            images_in_page.append(image_filename)

        #write image
        pages_images.append(images_in_page)


    pages_content = []
    for page_num in range(len(pdf_file)):
        pages_content.append({
            "slide_number": page_num + 1,
            "text": pages_text[page_num],
            "images": pages_images[page_num]
        })

    return pages_content







#create an input dir
inp_name = 'INPUTS'
folder_cr(inp_name,f' {inp_name} already exists!')

#create an output dir
out_name = 'OUTPUTS'
folder_cr(out_name,f' {out_name} already exists!')



list_to_do = os.listdir(path=inp_name)


for file in list_to_do:

    dot = file.find('.')
    middle_folder = file[:dot] + '_report'

    folder_cr('.' + '\\' + out_name + '\\' + middle_folder,f'Report folder for the {file[:dot]} already exists!')

    if '.pdf' in file:
        pdf = pymupdf.open(inp_name + '/' + file)
        a = pdf_converter(pdf,out_name,middle_folder)

    if '.pptx' in file:
        presentation = Presentation(inp_name + '/' + file)
        a = pptx_converter(presentation,out_name,middle_folder)

    with open(out_name + '/' + middle_folder + "/report.txt", "w",encoding="utf-8") as my_file:
        for each in a:
            my_file.write(str(each) + '\n')

